"""MarkItDown parser wrapper with a local text-extraction fallback."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Dict

SUPPORTED_FORMATS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    ".html": "html",
    ".htm": "html",
    ".txt": "text",
    ".md": "markdown",
}


class DocumentParseError(Exception):
    """Base error for document parsing failures."""


class UnsupportedFormatError(DocumentParseError):
    """Raised when the file suffix is not supported."""


class MarkItDownParser:
    """Convert local enterprise documents to normalized Markdown text."""

    def __init__(self) -> None:
        self._converter = None

    def convert(self, file_path: str) -> Dict[str, Any]:
        path = Path(file_path).expanduser()
        if not path.is_file():
            raise FileNotFoundError(f"document not found: {path}")

        suffix = path.suffix.lower()
        if suffix not in SUPPORTED_FORMATS:
            raise UnsupportedFormatError(
                f"unsupported document format: {suffix or 'unknown'}"
            )

        converter = self._get_converter()
        if converter is None:
            content = self._fallback_extract(path)
            engine = "fallback"
        else:
            try:
                result = converter.convert(str(path))
                content = getattr(result, "text_content", "") or ""
                engine = "markitdown"
            except Exception as exc:
                raise DocumentParseError(f"markitdown parse failed: {exc}") from exc

        stat = path.stat()
        return {
            "filename": path.name,
            "content": content.strip(),
            "format": SUPPORTED_FORMATS[suffix],
            "metadata": {
                "engine": engine,
                "path": str(path),
                "size": stat.st_size,
                "suffix": suffix,
            },
        }

    def _get_converter(self) -> Any:
        if self._converter is not None:
            return self._converter
        try:
            from markitdown import MarkItDown

            self._converter = MarkItDown()
            return self._converter
        except Exception:
            return None

    def _fallback_extract(self, path: Path) -> str:
        suffix = path.suffix.lower()
        if suffix == ".txt" or suffix == ".md":
            return path.read_text(encoding="utf-8", errors="replace")
        if suffix == ".docx":
            return self._extract_docx(path)
        if suffix == ".xlsx":
            return self._extract_xlsx(path)
        if suffix == ".pdf":
            return self._extract_pdf(path)
        if suffix in {".html", ".htm"}:
            return self._extract_html(path)
        if suffix == ".pptx":
            return self._extract_pptx(path)
        raise DocumentParseError(f"fallback parser is unavailable for {suffix}")

    @staticmethod
    def _extract_docx(path: Path) -> str:
        from docx import Document

        document = Document(str(path))
        lines = [paragraph.text.strip() for paragraph in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                lines.append(" | ".join(cell.text.strip() for cell in row.cells))
        return "\n".join(line for line in lines if line)

    @staticmethod
    def _extract_xlsx(path: Path) -> str:
        from openpyxl import load_workbook

        workbook = load_workbook(path, data_only=True, read_only=True)
        lines = []
        for worksheet in workbook.worksheets:
            lines.append(f"## {worksheet.title}")
            for row in worksheet.iter_rows(values_only=True):
                values = ["" if cell is None else str(cell).strip() for cell in row]
                if any(values):
                    lines.append(" | ".join(value for value in values if value))
        workbook.close()
        return "\n".join(lines)

    @staticmethod
    def _extract_pdf(path: Path) -> str:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages).strip()

    @staticmethod
    def _extract_html(path: Path) -> str:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="replace"), "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text("\n", strip=True)
        return re.sub(r"\n{3,}", "\n\n", text).strip()

    @staticmethod
    def _extract_pptx(path: Path) -> str:
        from pptx import Presentation

        presentation = Presentation(str(path))
        lines = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"## Slide {slide_index}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines)
